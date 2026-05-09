"""
Generate PowerPoint Presentation for BookWise AI Project
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BG = RGBColor(15, 23, 42)  # Dark navy
ACCENT_BLUE = RGBColor(59, 130, 246)  # Bright blue
ACCENT_PURPLE = RGBColor(147, 51, 234)  # Purple
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(226, 232, 240)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Add slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    
    for point in left_points:
        p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.2), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    
    for point in right_points:
        p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    return slide

# Slide 1: Title
add_title_slide(prs, "📚 BookWise AI", "Intelligent Book Recommender Chatbot\nPowered by Google Gemini 2.0 Flash")

# Slide 2: Problem Statement
add_content_slide(prs, "What Problem Does It Solve?", [
    "😕 Users struggle to find the right book from thousands of choices",
    "⏳ Manual searching takes time and often gives generic suggestions",
    "📌 Existing tools rarely understand mood, context, or reading level",
    "🌍 Many readers want recommendations in their own language",
    "📄 Users also need a way to ask questions about uploaded content"
])

# Slide 3: What We Achieved
add_content_slide(prs, "What We Achieved", [
    "🤖 Built an AI chatbot that gives personalized book recommendations",
    "🧠 Added session-based memory so the chatbot remembers context",
    "📚 Connected a local book knowledge base for reliable suggestions",
    "🌐 Enabled multilingual support with automatic language detection",
    "📄 Added document Q&A so users can ask questions from uploaded text",
    "🎛️ Made response creativity adjustable with temperature and top-p sliders"
])

# Slide 4: How It Works
add_two_column_slide(prs, "How It Works", 
    "User Side", 
    [
        "Type a question in the web interface",
        "Adjust temperature and top-p if needed",
        "Upload a document for Q&A",
        "Receive personalized book suggestions"
    ],
    "AI Side",
    [
        "Flask receives the request",
        "Context is built from history + knowledge base",
        "Gemini 2.0 Flash generates the response",
        "Response is sent back to the chat UI"
    ]
)

# Slide 5: Why AI Was Needed
add_content_slide(prs, "Why AI Was Needed", [
    "🤔 Traditional search cannot understand mood, intent, or context well",
    "📚 Readers want personalized recommendations instead of generic lists",
    "⏱️ AI makes book discovery faster and more interactive",
    "🌍 Users expect support for different languages and reading preferences",
    "📄 The project also needed a smart way to answer questions from documents"
])

# Slide 6: AI Pipeline / Architecture
add_two_column_slide(prs, "AI Pipeline and Architecture",
    "Input and Context",
    [
        "User message comes from the web chat",
        "Session history is collected",
        "Book knowledge base is loaded",
        "Uploaded document text is added if available"
    ],
    "Model and Output",
    [
        "Gemini 2.0 Flash receives the full prompt",
        "Temperature and top-p shape the answer",
        "The model generates a personalized reply",
        "The response is rendered back in the browser"
    ]
)

# Slide 7: Core AI / ML Part
add_two_column_slide(prs, "AI and Machine Learning", 
    "AI Model Used",
    [
        "Gemini 2.0 Flash powers the chatbot responses",
        "Google Generative AI API connects the app to the model",
        "Temperature controls creativity and response variety",
        "Top-P controls how focused or diverse the output is",
        "Max output tokens keep responses detailed but controlled"
    ],
    "Machine Learning & NLP",
    [
        "Natural language processing understands user messages",
        "Intent recognition identifies book, mood, or author requests",
        "Context-aware reasoning uses chat history for better replies",
        "Language detection lets the chatbot answer in the same language",
        "Prompt engineering shapes responses with book-specific rules"
    ]
)

# Slide 8: Key Features
add_content_slide(prs, "Key Features", [
    "📖 Book recommendations by genre, mood, author, and reading style",
    "📚 Spoiler-free summaries with reasoning for each suggestion",
    "🧠 Conversation memory for better follow-up answers",
    "🌐 Multilingual responses for broader accessibility",
    "📄 Document Q&A for uploaded files",
    "💾 Local JSON knowledge base for stable reference data"
])

# Slide 9: Technology Stack
add_two_column_slide(prs, "Technology Stack", 
    "Frontend",
    [
        "HTML5, CSS3, Vanilla JavaScript",
        "Dark theme with glassmorphism",
        "Responsive chat interface"
    ],
    "Backend and AI",
    [
        "Python 3 and Flask",
        "Google Gemini 2.0 Flash",
        "LangDetect, Flask-CORS, JSON knowledge base"
    ]
)

# Slide 10: Impact and Conclusion
add_content_slide(prs, "Impact", [
    "✅ Saves time by turning book discovery into a conversation",
    "✅ Gives better recommendations by using context and preferences",
    "✅ Helps users discover books that match mood, language, and level",
    "✅ Works as a practical AI assistant for readers",
    "✅ Demonstrates a complete hybrid AI + web application project"
])

# Slide 11: Closing
add_title_slide(prs, "Thank You!", "GitHub: github.com/Pahar12/Book-recommender-chatbot\nRun locally at: http://localhost:5001")

# Save presentation
output_path = "BookWise_AI_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created: {output_path}")
